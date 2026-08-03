import io
import re
import warnings
import time
import requests
from datetime import datetime
from zoneinfo import ZoneInfo

import joblib
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
import streamlit as st
from PIL import Image
from matplotlib.backends.backend_pdf import PdfPages
from scipy.signal import peak_widths
from scipy.stats import entropy, kurtosis, skew
from sklearn.base import clone
from sklearn.cross_decomposition import PLSRegression
from sklearn.decomposition import PCA
from sklearn.ensemble import (ExtraTreesRegressor, GradientBoostingRegressor,
                              HistGradientBoostingRegressor, RandomForestRegressor)
from sklearn.feature_selection import f_regression, mutual_info_regression
from sklearn.impute import SimpleImputer
from sklearn.inspection import permutation_importance
from sklearn.linear_model import ElasticNet, Lasso, LinearRegression, Ridge
from sklearn.metrics import (accuracy_score, balanced_accuracy_score,
                             confusion_matrix, mean_absolute_error,
                             mean_squared_error, r2_score, roc_curve, auc)
from sklearn.model_selection import train_test_split
from sklearn.neighbors import KNeighborsRegressor
from sklearn.pipeline import make_pipeline
from sklearn.preprocessing import MinMaxScaler, StandardScaler
from sklearn.svm import SVR

warnings.filterwarnings("ignore")
RANDOM_STATE = 42
META = ["sample_id", "pH", "temperature_C", "replicate", "source_file"]

st.set_page_config(page_title="CV pH Prediction", page_icon="⚗️", layout="wide")
st.markdown("""
<style>
.stApp {background:radial-gradient(circle at 8% 5%,rgba(0,189,214,.14),transparent 24%),radial-gradient(circle at 92% 8%,rgba(126,61,181,.15),transparent 25%),linear-gradient(135deg,#f3fbff 0%,#fff8f5 50%,#f6f0ff 100%);}
[data-testid="stHeader"] {background: rgba(255,255,255,.72);}
.hero {padding: 1.45rem 1.6rem; border-radius: 22px; color: white;
 background: linear-gradient(115deg,#073b66,#087ea4 48%,#6f42c1);
 box-shadow: 0 12px 30px rgba(15,75,120,.22); margin-bottom: 1rem;position:relative;overflow:hidden;animation:heroGlow 5s ease-in-out infinite;}
@keyframes heroGlow{0%,100%{box-shadow:0 12px 30px rgba(15,75,120,.22)}50%{box-shadow:0 16px 42px rgba(111,66,193,.38)}}
.hero h1 {font-size: 2rem; margin:0 0 .25rem 0; color:white;}
.hero p {margin:0; opacity:.92; font-size:1.02rem;}
.info-card {border-radius:18px; padding:1rem 1.15rem; color:white; min-height:118px;
 box-shadow:0 8px 22px rgba(35,55,80,.16);}
.blue-card {background:linear-gradient(135deg,#0061a8,#00a6c7);}
.purple-card {background:linear-gradient(135deg,#6338a8,#b149b9);}
.green-card {background:linear-gradient(135deg,#047857,#12a884);}
.orange-card {background:linear-gradient(135deg,#d35a20,#f59e0b);}
.prediction-card {padding:1.3rem; border-radius:20px; text-align:center; color:white;
 background:linear-gradient(135deg,#5928a0,#9b2fae,#d94878); box-shadow:0 10px 28px rgba(92,40,150,.25);}
.prediction-card .value {font-size:3rem; font-weight:800; line-height:1.1;}
.hero-badges{display:flex;gap:.55rem;flex-wrap:wrap;margin-top:1rem}.hero-badge{padding:.35rem .72rem;border-radius:999px;background:rgba(255,255,255,.16);border:1px solid rgba(255,255,255,.27);font-size:.82rem;font-weight:650;backdrop-filter:blur(8px)}
.status-dot{display:inline-block;width:9px;height:9px;border-radius:50%;background:#28f49a;margin-right:6px;box-shadow:0 0 0 rgba(40,244,154,.55);animation:statusPulse 1.5s infinite}@keyframes statusPulse{0%{box-shadow:0 0 0 0 rgba(40,244,154,.55)}70%{box-shadow:0 0 0 10px rgba(40,244,154,0)}100%{box-shadow:0 0 0 0 rgba(40,244,154,0)}}
.section-card{background:rgba(255,255,255,.84);border:1px solid rgba(255,255,255,.92);border-radius:18px;padding:.85rem 1rem;box-shadow:0 8px 25px rgba(30,62,95,.09);backdrop-filter:blur(10px);margin:.35rem 0 .8rem}
div[data-testid="stMetric"] {background:white; border:1px solid #dfe9f2; padding:14px;
 border-radius:16px; box-shadow:0 5px 15px rgba(20,60,90,.08);}
.stButton>button {border:0;border-radius:13px;font-weight:750;color:white;background:linear-gradient(110deg,#087ea4,#6f42c1,#d33b83);background-size:200% 200%;box-shadow:0 7px 18px rgba(72,64,150,.22);transition:.25s;animation:buttonFlow 4s ease infinite;}
.stButton>button:hover {transform:translateY(-2px);box-shadow:0 10px 25px rgba(72,64,150,.34);color:white;}
@keyframes buttonFlow{0%,100%{background-position:0% 50%}50%{background-position:100% 50%}}
[data-testid="stTabs"] button[aria-selected="true"]{color:#6f42c1;font-weight:800;border-bottom:3px solid #d33b83;}
html,body,[class*="css"]{font-family:Inter,"Segoe UI",Roboto,Arial,sans-serif;color:#17324d;}
.block-container{max-width:1450px;padding-top:1.25rem;padding-bottom:3rem;}
h1,h2,h3{letter-spacing:-.025em;color:#17324d}h2{margin-top:1.25rem!important}h3{padding-top:.35rem}
[data-testid="stSidebar"]{background:linear-gradient(180deg,#062f58 0%,#075879 48%,#44327a 100%);border-right:1px solid rgba(255,255,255,.18);}
[data-testid="stSidebar"] *{color:#f7fbff!important}[data-testid="stSidebar"] input{color:#17324d!important;background:white!important}
[data-testid="stSidebar"] [data-baseweb="select"] *{color:#17324d!important}
[data-testid="stFileUploader"]{background:rgba(255,255,255,.83);border:2px dashed #60a9ce;border-radius:18px;padding:.55rem;box-shadow:0 7px 20px rgba(29,78,110,.08);}
[data-testid="stFileUploader"]:hover{border-color:#8b53c6;box-shadow:0 10px 28px rgba(111,66,193,.14)}
[data-baseweb="select"]>div,[data-testid="stNumberInput"] input,[data-testid="stTextInput"] input{background:rgba(255,255,255,.92)!important;border-color:#cdddeb!important;border-radius:12px!important;box-shadow:0 3px 10px rgba(30,70,100,.05)}
[data-baseweb="select"]>div:focus-within,[data-testid="stNumberInput"] input:focus,[data-testid="stTextInput"] input:focus{border-color:#7952b3!important;box-shadow:0 0 0 3px rgba(121,82,179,.13)!important}
[data-testid="stTabs"]{background:rgba(255,255,255,.67);border:1px solid rgba(255,255,255,.9);border-radius:18px;padding:.25rem .7rem .7rem;box-shadow:0 7px 24px rgba(30,65,95,.07)}
[data-testid="stTabs"] button{font-weight:650;padding:.8rem 1rem;transition:.2s}[data-testid="stTabs"] button:hover{color:#d33b83;transform:translateY(-1px)}
[data-testid="stMetric"] label{color:#557086!important;font-weight:650}[data-testid="stMetricValue"]{color:#173f70!important;font-weight:800}
[data-testid="stDataFrame"]{border:1px solid #dbe7f0;border-radius:14px;overflow:hidden;box-shadow:0 6px 18px rgba(25,65,95,.07)}
[data-testid="stPlotlyChart"], [data-testid="stImage"], [data-testid="stPyplotGlobalUse"]{background:rgba(255,255,255,.88);border:1px solid rgba(225,235,244,.95);border-radius:18px;padding:.35rem;box-shadow:0 8px 24px rgba(25,65,95,.08)}
[data-testid="stAlert"]{border-radius:14px;border:0;box-shadow:0 5px 16px rgba(30,70,100,.07)}
.stDownloadButton>button{background:linear-gradient(110deg,#087ea4,#047857)!important;animation:none!important}
hr{border:0;height:1px;background:linear-gradient(90deg,transparent,#9cbfd5,transparent);margin:1.4rem 0}
@media(max-width:768px){.hero{padding:1.2rem}.hero h1{font-size:1.55rem}.hero-badges{gap:.35rem}.hero-badge{font-size:.72rem}.block-container{padding-left:.7rem;padding-right:.7rem}}
</style>
<div class="hero">
  <h1>⚗️ AI-TC pH Sensor App</h1>
  <p>AI-powered temperature-compensated electrochemical pH prediction</p>
  <div class="hero-badges"><span class="hero-badge"><span class="status-dot"></span>Analysis engine ready</span><span class="hero-badge">🌡️ Temperature compensated</span><span class="hero-badge">🧠 Explainable ML</span><span class="hero-badge">📍 Geo-tagged reports</span></div>
</div>
""", unsafe_allow_html=True)
st.caption("Upload → feature extraction → intelligent selection → model training → smart prediction")
st.markdown("""
<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(180px,1fr));gap:.75rem;margin:.65rem 0 1.05rem">
  <div class="section-card"><div style="font-size:1.45rem">📈</div><b>Electrochemistry</b><br><small style="color:#60798c">Raw and mean CV analytics</small></div>
  <div class="section-card"><div style="font-size:1.45rem">🧬</div><b>Smart Features</b><br><small style="color:#60798c">3D feature exploration</small></div>
  <div class="section-card"><div style="font-size:1.45rem">🤖</div><b>Model Studio</b><br><small style="color:#60798c">Training and comparison</small></div>
  <div class="section-card"><div style="font-size:1.45rem">🎯</div><b>Prediction Lab</b><br><small style="color:#60798c">Live stream and reports</small></div>
</div>
""",unsafe_allow_html=True)


def ph_from_name(name):
    match = re.search(r"ph[\s_\-]*(\d+(?:\.\d+)?)", name, re.I)
    if not match:
        raise ValueError(f"Cannot identify pH in '{name}'. Rename it like pH 3.xlsx.")
    return float(match.group(1))


def download_onedrive_file(shared_url, filename):
    """Download an anonymously shared OneDrive file into a file-like object."""
    shared_url=shared_url.strip()
    if not shared_url.startswith(("https://","http://")):
        raise ValueError("The OneDrive link must start with https://")
    separator="&" if "?" in shared_url else "?"
    download_url=shared_url if "download=1" in shared_url else shared_url+separator+"download=1"
    response=requests.get(download_url,timeout=45,allow_redirects=True)
    response.raise_for_status()
    content_type=response.headers.get("content-type","").lower()
    if "text/html" in content_type:
        raise ValueError("OneDrive returned a webpage instead of the file. Share it as 'Anyone with the link can view'.")
    file_object=io.BytesIO(response.content); file_object.name=filename.strip(); return file_object


def read_cv_file(uploaded):
    ph = ph_from_name(uploaded.name)
    raw = pd.read_excel(uploaded, sheet_name=0)
    records, counters, columns = [], {}, list(raw.columns)
    for j, col in enumerate(columns):
        match = re.match(r"^I\s*(4|20|21|30|40|50|60|70|80)\b", str(col).strip(), re.I)
        if not match or j == 0:
            continue
        temperature = int(match.group(1))
        v = pd.to_numeric(raw[columns[j - 1]], errors="coerce")
        i = pd.to_numeric(raw[col], errors="coerce")
        valid = v.notna() & i.notna()
        v, i = v[valid].to_numpy(float), i[valid].to_numpy(float)
        if len(v) < 10:
            continue
        counters[temperature] = counters.get(temperature, 0) + 1
        rep = counters[temperature]
        records.append({"sample_id": f"pH{ph:g}_T{temperature}_R{rep}", "pH": ph,
                        "temperature_C": temperature, "replicate": rep,
                        "potential_V": v, "current_A": i, "source_file": uploaded.name})
    if not records:
        raise ValueError(f"No current columns found in '{uploaded.name}'. Expected I4, I20, I30, etc.")
    return records


def split_group(group):
    group = group.sample(frac=1, random_state=RANDOM_STATE).copy()
    n = len(group)
    if n < 3:
        group["dataset"] = "Training"
        return group
    n_test = max(1, round(n * .15)); n_val = max(1, round(n * .15))
    while n_test + n_val >= n:
        if n_val > 1: n_val -= 1
        elif n_test > 1: n_test -= 1
        else: break
    train_val, test = train_test_split(group, test_size=n_test, random_state=RANDOM_STATE)
    train, val = train_test_split(train_val, test_size=n_val, random_state=RANDOM_STATE)
    train["dataset"], val["dataset"], test["dataset"] = "Training", "Validation", "Testing"
    return pd.concat([train, val, test])


def extract_features(row):
    v = np.asarray(row["potential_V"], float); i = np.asarray(row["current_A"], float) * 1e9
    valid = np.isfinite(v) & np.isfinite(i); v, i = v[valid], i[valid]
    if len(v) < 10: raise ValueError(f"Too few points in {row['sample_id']}")
    ox, red = int(np.argmax(i)), int(np.argmin(i)); ipa, ipc, epa, epc = i[ox], i[red], v[ox], v[red]
    with np.errstate(all="ignore"):
        didv = np.gradient(i, v); d2 = np.gradient(didv, v)
    def width(signal, index):
        try:
            left, right = peak_widths(signal, [index], rel_height=.5)[2:4]
            x = np.arange(len(v)); return abs(np.interp(right[0], x, v) - np.interp(left[0], x, v))
        except Exception: return np.nan
    ox_area = abs(np.trapezoid(np.clip(i, 0, None), v)); red_area = abs(np.trapezoid(np.clip(-i, 0, None), v))
    turn, loop_area = int(np.argmax(v)), np.nan
    if 2 < turn < len(v) - 3:
        vf, iff = v[:turn + 1], i[:turn + 1]
        vr, indices = np.unique(v[turn:][::-1], return_index=True); irr = i[turn:][::-1][indices]
        lo, hi = max(vf.min(), vr.min()), min(vf.max(), vr.max())
        if lo < hi:
            grid = np.linspace(lo, hi, 300)
            loop_area = np.trapezoid(abs(np.interp(grid, vf, iff) - np.interp(grid, vr, irr)), grid)
    probability = (abs(i) + 1e-12); probability /= probability.sum()
    out = {"sample_id": row["sample_id"], "pH": row["pH"], "temperature_C": row["temperature_C"],
           "replicate": row["replicate"], "source_file": row["source_file"],
           "oxidation_peak_current_nA": ipa, "reduction_peak_current_nA": ipc,
           "oxidation_peak_voltage_V": epa, "reduction_peak_voltage_V": epc,
           "peak_separation_V": epa-epc, "peak_current_difference_nA": ipa-ipc,
           "peak_current_ratio": ipa/abs(ipc) if ipc else np.nan,
           "oxidation_FWHM_V": width(i, ox), "reduction_FWHM_V": width(-i, red),
           "oxidation_sharpness_nA_per_V2": abs(d2[ox]), "reduction_sharpness_nA_per_V2": abs(d2[red]),
           "oxidation_area_nA_V": ox_area, "reduction_area_nA_V": red_area,
           "total_peak_area_nA_V": ox_area+red_area, "CV_loop_area_nA_V": loop_area,
           "absolute_signal_area": np.trapezoid(abs(i), x=np.arange(len(i))),
           "maximum_current_nA": np.max(i), "minimum_current_nA": np.min(i), "mean_current_nA": np.mean(i),
           "median_current_nA": np.median(i), "current_range_nA": np.ptp(i), "current_std_nA": np.std(i, ddof=1),
           "current_variance_nA2": np.var(i, ddof=1), "current_RMS_nA": np.sqrt(np.mean(i**2)),
           "current_skewness": skew(i), "current_kurtosis": kurtosis(i),
           "maximum_slope_nA_per_V": np.nanmax(didv), "minimum_slope_nA_per_V": np.nanmin(didv),
           "mean_slope_nA_per_V": np.nanmean(didv), "slope_std_nA_per_V": np.nanstd(didv),
           "maximum_curvature": np.nanmax(abs(d2)), "mean_curvature": np.nanmean(abs(d2)),
           "voltage_range_V": np.ptp(v), "number_of_points": len(v),
           "zero_crossings": np.sum(np.diff(np.signbit(i)) != 0), "signal_energy_nA2": np.sum(i**2),
           "signal_entropy": entropy(probability),
           "curve_length": np.sum(np.sqrt(np.diff(v)**2 + np.diff(i)**2)), "roughness_nA": np.std(np.diff(i))}
    return out


def rank_features(train_features):
    candidates = [c for c in train_features.columns if c not in META and pd.api.types.is_numeric_dtype(train_features[c])]
    x = train_features[candidates].replace([np.inf, -np.inf], np.nan)
    x = x.drop(columns=[c for c in x if x[c].nunique(dropna=True) <= 1])
    clean = pd.DataFrame(SimpleImputer(strategy="median").fit_transform(x), columns=x.columns, index=x.index)
    y = train_features["pH"]
    fscore, pvalue = f_regression(clean, y)
    mi = mutual_info_regression(clean, y, random_state=RANDOM_STATE)
    forest = RandomForestRegressor(n_estimators=300, random_state=RANDOM_STATE, n_jobs=-1).fit(clean, y)
    pi = permutation_importance(forest, clean, y, n_repeats=10, random_state=RANDOM_STATE, n_jobs=-1)
    ranking = pd.DataFrame({"Feature": clean.columns, "F_Score": fscore, "P_Value": pvalue,
                            "Adjusted_P_Value": np.minimum(pvalue*clean.shape[1], 1),
                            "Mutual_Information": mi, "RF_Importance": forest.feature_importances_,
                            "Permutation_Importance": pi.importances_mean})
    scores = ["F_Score", "Mutual_Information", "RF_Importance", "Permutation_Importance"]
    scaled = MinMaxScaler().fit_transform(ranking[scores].fillna(0))
    ranking["Combined_Score"] = scaled.mean(axis=1)
    ranking["Statistically_Significant"] = ranking["Adjusted_P_Value"] < .05
    ranking = ranking.sort_values("Combined_Score", ascending=False).reset_index(drop=True)
    ranking.insert(0, "Rank", np.arange(1, len(ranking)+1))
    return ranking, clean


def model_set(feature_count, train_count):
    scaled = lambda model: make_pipeline(SimpleImputer(strategy="median"), StandardScaler(), model)
    plain = lambda model: make_pipeline(SimpleImputer(strategy="median"), model)
    return {"Linear Regression": scaled(LinearRegression()), "Ridge Regression": scaled(Ridge(alpha=1)),
            "LASSO Regression": scaled(Lasso(alpha=.01, max_iter=20000)),
            "Elastic Net": scaled(ElasticNet(alpha=.01, l1_ratio=.5, max_iter=20000)),
            "PLS Regression": scaled(PLSRegression(n_components=max(1, min(5, feature_count, train_count-1)))),
            "SVR": scaled(SVR(kernel="rbf", C=10, epsilon=.1)),
            "KNN Regression": scaled(KNeighborsRegressor(n_neighbors=max(1, min(5, train_count)), weights="distance")),
            "Random Forest": plain(RandomForestRegressor(n_estimators=300, random_state=RANDOM_STATE, n_jobs=-1)),
            "Extra Trees": plain(ExtraTreesRegressor(n_estimators=300, random_state=RANDOM_STATE, n_jobs=-1)),
            "Gradient Boosting": plain(GradientBoostingRegressor(n_estimators=300, learning_rate=.05, max_depth=3, random_state=RANDOM_STATE)),
            "Histogram Gradient Boosting": plain(HistGradientBoostingRegressor(max_iter=300, learning_rate=.05, random_state=RANDOM_STATE))}


def nearest(values, classes):
    values = np.asarray(values).reshape(-1); classes = np.asarray(classes)
    return classes[np.abs(values[:, None]-classes[None, :]).argmin(axis=1)]


def mean_cv(curves, points=400):
    """Mean CV on normalized scan progress, preserving forward/reverse branches."""
    grid=np.linspace(0,1,points); voltages=[]; currents=[]
    for _,row in curves.iterrows():
        voltage=np.asarray(row["potential_V"],float); current=np.asarray(row["current_A"],float)*1e9
        valid=np.isfinite(voltage)&np.isfinite(current); voltage,current=voltage[valid],current[valid]
        if len(voltage)<10: continue
        progress=np.linspace(0,1,len(voltage)); voltages.append(np.interp(grid,progress,voltage)); currents.append(np.interp(grid,progress,current))
    if not currents: return None
    voltages=np.vstack(voltages); currents=np.vstack(currents)
    return voltages.mean(0),currents.mean(0),currents.std(0,ddof=1) if len(currents)>1 else np.zeros(points)


def prediction_gif(raw,predicted_ph,actual_ph):
    voltage=np.asarray(raw["potential_V"],float); current=np.asarray(raw["current_A"],float)*1e9; frames=[]
    reference=predicted_ph-.6 if actual_ph is None else actual_ph; start=(reference+predicted_ph)/2
    for fraction in np.linspace(.08,1,18):
        end=max(3,int(len(voltage)*fraction)); evolving=start+(predicted_ph-start)*fraction
        fig,(ax1,ax2)=plt.subplots(1,2,figsize=(9,3.6),gridspec_kw={"width_ratios":[1.7,1]})
        ax1.plot(voltage[:end],current[:end],color="#6f2dbd",lw=2.2); ax1.set_xlim(voltage.min(),voltage.max()); pad=max(np.ptp(current)*.08,.1); ax1.set_ylim(current.min()-pad,current.max()+pad)
        ax1.set(xlabel="Potential (V)",ylabel="Current (nA)",title="Live CV acquisition"); ax1.grid(alpha=.2)
        ax2.scatter(.5,evolving,s=600,c="#d33b83",edgecolor="white",lw=3); ax2.text(.5,evolving,f"{evolving:.2f}",ha="center",va="center",color="white",weight="bold")
        if actual_ph is not None: ax2.axhline(actual_ph,color="#00856a",ls="--",label=f"Actual {actual_ph:g}"); ax2.legend(fontsize=8)
        ax2.set_xlim(0,1); ax2.set_ylim(min(reference,predicted_ph)-1,max(reference,predicted_ph)+1); ax2.set_xticks([]); ax2.set_ylabel("pH"); ax2.set_title("Evolving ML prediction")
        fig.suptitle("AI-TC pH Sensor App",weight="bold"); fig.tight_layout(); buffer=io.BytesIO(); fig.savefig(buffer,format="png",dpi=85,bbox_inches="tight"); plt.close(fig); buffer.seek(0); frames.append(Image.open(buffer).convert("P"))
    output=io.BytesIO(); frames[0].save(output,format="GIF",save_all=True,append_images=frames[1:],duration=160,loop=0); return output.getvalue()


def create_pdf_and_powerpoint(location, timestamp, results, test, predicted, model_name, include_powerpoint=True):
    """Create a compact analysis report only when the user requests it."""
    figures=[]
    cover,ax=plt.subplots(figsize=(11,7)); ax.axis("off")
    ax.text(.05,.86,"AI-TC pH Sensor App",fontsize=28,fontweight="bold",color="#153b66")
    ax.text(.05,.74,"Temperature-compensated electrochemical pH prediction",fontsize=16,color="#6f2dbd")
    ax.text(.05,.56,f"Location: {location}\nDate and time: {timestamp}\nSelected model: {model_name}",fontsize=14,linespacing=1.7)
    best=results.iloc[0]
    ax.text(.05,.30,f"Best validation model: {best['Model']}\nValidation RMSE: {best['Validation_RMSE']:.4f}\nTest RMSE: {best['Test_RMSE']:.4f}",fontsize=15,linespacing=1.5,bbox=dict(boxstyle="round,pad=.7",facecolor="#eef7ff",edgecolor="#7ab8df"))
    figures.append(("Analysis summary",cover))

    comparison,ax=plt.subplots(figsize=(10,6)); ordered=results.sort_values("Test_RMSE",ascending=True)
    ax.barh(ordered.Model,ordered.Test_RMSE,color=plt.cm.viridis(np.linspace(.15,.9,len(ordered))))
    ax.set_xlabel("Test RMSE (pH)"); ax.set_title("Regression Model Comparison",fontweight="bold"); ax.grid(axis="x",alpha=.2); comparison.tight_layout()
    figures.append(("Model comparison",comparison))

    prediction,ax=plt.subplots(figsize=(8,6)); actual=test.pH.to_numpy(float); predicted=np.asarray(predicted,float)
    scatter=ax.scatter(actual,predicted,c=test.temperature_C,cmap="turbo",s=75,edgecolor="white",linewidth=.6)
    low=min(actual.min(),predicted.min()); high=max(actual.max(),predicted.max()); ax.plot([low,high],[low,high],"k--")
    ax.set_xlabel("Actual pH"); ax.set_ylabel("Predicted pH"); ax.set_title(f"{model_name}: Actual vs Predicted",fontweight="bold"); prediction.colorbar(scatter,ax=ax,label="Temperature (°C)"); ax.grid(alpha=.2); prediction.tight_layout()
    figures.append(("Prediction performance",prediction))

    pdf=io.BytesIO()
    with PdfPages(pdf) as pages:
        for _,figure in figures: pages.savefig(figure,bbox_inches="tight")

    if not include_powerpoint:
        for _,figure in figures: plt.close(figure)
        return pdf.getvalue(),None

    # Imported here so a temporary package issue does not stop the whole app.
    from pptx import Presentation
    from pptx.util import Inches, Pt
    deck=Presentation(); deck.slide_width=Inches(13.333); deck.slide_height=Inches(7.5)
    for title,figure in figures:
        image=io.BytesIO(); figure.savefig(image,format="png",dpi=170,bbox_inches="tight",facecolor="white"); image.seek(0)
        slide=deck.slides.add_slide(deck.slide_layouts[6]); box=slide.shapes.add_textbox(Inches(.5),Inches(.12),Inches(12),Inches(.55)); paragraph=box.text_frame.paragraphs[0]; paragraph.text=title; paragraph.font.size=Pt(24); paragraph.font.bold=True
        slide.shapes.add_picture(image,Inches(.65),Inches(.75),width=Inches(12),height=Inches(6.35))
    powerpoint=io.BytesIO(); deck.save(powerpoint)
    for _,figure in figures: plt.close(figure)
    return pdf.getvalue(),powerpoint.getvalue()


def create_live_reports(location, timestamp, model_name, predicted_ph, temperature, raw, features, actual_ph=None):
    """Create self-contained PDF and PowerPoint reports for an external MIT App CSV."""
    actual_text="Not supplied" if actual_ph is None else f"{actual_ph:.3f}"
    error_text="Not available" if actual_ph is None else f"{abs(predicted_ph-actual_ph):.3f}"
    figures=[]
    cover,ax=plt.subplots(figsize=(11,7)); ax.axis("off")
    ax.text(.05,.86,"AI-TC pH Sensor App",fontsize=28,fontweight="bold",color="#153b66")
    ax.text(.05,.74,"Live MIT App / OneDrive test prediction",fontsize=17,color="#6f2dbd")
    ax.text(.05,.56,f"Location: {location}\nDate and time: {timestamp}\nModel: {model_name}\nTemperature: {temperature:.2f} °C",fontsize=14,linespacing=1.55)
    ax.text(.05,.26,f"Predicted pH: {predicted_ph:.3f}\nActual pH: {actual_text}\nAbsolute error: {error_text}",fontsize=17,linespacing=1.5,bbox=dict(boxstyle="round,pad=.7",facecolor="#eef7ff",edgecolor="#7ab8df"))
    figures.append(("Live prediction summary",cover))
    curve,ax=plt.subplots(figsize=(10,6)); ax.plot(raw["potential_V"],np.asarray(raw["current_A"])*1e9,color="#6f2dbd",lw=2.3); ax.fill_between(raw["potential_V"],np.asarray(raw["current_A"])*1e9,color="#00a6c7",alpha=.14)
    ax.set(xlabel="Potential (V)",ylabel="Current (nA)",title="Uploaded live-test CV curve"); ax.grid(alpha=.2); curve.tight_layout(); figures.append(("Raw live-test CV",curve))
    pdf=io.BytesIO()
    with PdfPages(pdf) as pages:
        for _,figure in figures: pages.savefig(figure,bbox_inches="tight")
    from pptx import Presentation
    from pptx.util import Inches,Pt
    deck=Presentation(); deck.slide_width=Inches(13.333); deck.slide_height=Inches(7.5)
    for title,figure in figures:
        picture=io.BytesIO(); figure.savefig(picture,format="png",dpi=170,bbox_inches="tight",facecolor="white"); picture.seek(0)
        slide=deck.slides.add_slide(deck.slide_layouts[6]); box=slide.shapes.add_textbox(Inches(.5),Inches(.12),Inches(12),Inches(.55)); p=box.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(24); p.font.bold=True
        slide.shapes.add_picture(picture,Inches(.65),Inches(.75),width=Inches(12),height=Inches(6.35))
    powerpoint=io.BytesIO(); deck.save(powerpoint)
    for _,figure in figures: plt.close(figure)
    return pdf.getvalue(),powerpoint.getvalue()


def excel_bytes(sheets):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        for name, data in sheets.items(): data.to_excel(writer, sheet_name=name[:31], index=False)
    return output.getvalue()


with st.sidebar:
    st.header("📍 Measurement context")
    location_name = st.text_input("Location name", "Cork, Ireland")
    latitude = st.number_input("Latitude", value=51.8985, format="%.5f")
    longitude = st.number_input("Longitude", value=-8.4756, format="%.5f")
    timezone_name = st.selectbox(
        "Time zone",
        ["Europe/Dublin", "Asia/Kolkata", "Europe/London", "UTC"],
        index=0,
    )
    try:
        local_now = datetime.now(ZoneInfo(timezone_name))
    except Exception:
        local_now = datetime.now(ZoneInfo("UTC"))
    st.markdown(
        f"""<div class="info-card blue-card">
        <div style="font-size:1.8rem">📍 🕒</div>
        <b>{location_name}</b><br>
        {local_now.strftime('%d %B %Y')}<br>
        <span style="font-size:1.25rem;font-weight:750">{local_now.strftime('%I:%M:%S %p')}</span>
        </div>""",
        unsafe_allow_html=True,
    )
    st.map(pd.DataFrame({"lat": [latitude], "lon": [longitude]}), zoom=9)


upload_tab, feature_tab, selection_tab, training_tab, comparison_tab, prediction_tab, live_test_tab = st.tabs([
    "1 · Upload", "2 · Extract", "3 · Select", "4 · Train", "5 · Compare", "6 · Smart prediction", "7 · Live test data"])

with upload_tab:
    training_source=st.radio("Training-data source",["Device / OneDrive file picker","OneDrive shared links"],horizontal=True)
    uploads=[]
    if training_source=="Device / OneDrive file picker":
        uploads = st.file_uploader("Upload one Excel file for each pH", type=["xlsx", "xls"], accept_multiple_files=True,
                                   help="On a phone or Windows file picker, choose OneDrive and select the files. Filenames must contain pH, for example pH 3.xlsx.")
    else:
        st.info("Use one line per file in the format: pH 3.xlsx | https://your-OneDrive-link")
        training_links=st.text_area("OneDrive training-file links",placeholder="pH 3.xlsx | https://1drv.ms/...\npH 4.xlsx | https://1drv.ms/...")
        if st.button("Fetch Excel files from OneDrive",use_container_width=True):
            try:
                downloaded=[]
                for line in [item.strip() for item in training_links.splitlines() if item.strip()]:
                    if "|" not in line: raise ValueError("Each line must contain filename | OneDrive link")
                    filename,url=line.split("|",1); file_object=download_onedrive_file(url,filename)
                    downloaded.append({"name":filename.strip(),"bytes":file_object.getvalue()})
                if not downloaded: raise ValueError("Enter at least one OneDrive link.")
                st.session_state.onedrive_training_files=downloaded; st.success(f"Fetched {len(downloaded)} Excel files.")
            except Exception as exc: st.error(f"OneDrive download failed: {exc}")
        for stored in st.session_state.get("onedrive_training_files",[]):
            file_object=io.BytesIO(stored["bytes"]); file_object.name=stored["name"]; uploads.append(file_object)
    if st.button("Read uploaded CV data", type="primary", disabled=not uploads):
        try:
            records = []
            for uploaded in uploads: records.extend(read_cv_file(uploaded))
            data = pd.DataFrame(records)
            # Use an explicit loop instead of groupby.apply(). Pandas 3.x may
            # exclude grouping columns from the DataFrame passed to apply(),
            # which removes pH/temperature_C on Streamlit Cloud.
            split_parts = []
            for (ph_value, temperature), group in data.groupby(
                ["pH", "temperature_C"], sort=True
            ):
                group = group.copy()
                group["pH"] = ph_value
                group["temperature_C"] = temperature
                split_parts.append(split_group(group))
            split = pd.concat(split_parts, ignore_index=True)
            st.session_state.update(cv_data=data, split_data=split)
            for key in ["features", "ranking", "selected", "results", "models", "predictions"]: st.session_state.pop(key, None)
            st.success(f"Detected {len(data)} complete CV curves from {len(uploads)} files.")
        except Exception as exc: st.error(str(exc))
    if "cv_data" in st.session_state:
        data, split = st.session_state.cv_data, st.session_state.split_data
        a,b,c,d = st.columns(4)
        a.metric("CV curves", len(data)); b.metric("pH values", data.pH.nunique()); c.metric("Temperatures", data.temperature_C.nunique()); d.metric("Files", data.source_file.nunique())
        st.dataframe(pd.crosstab([data.pH], data.temperature_C), use_container_width=True)
        # Build the count table explicitly. This avoids an ambiguity in newer
        # pandas/seaborn versions when grouped data retains pH as an index level.
        split_counts = (
            split.reset_index(drop=True)
            .groupby(["pH", "dataset"], as_index=False)
            .size()
            .pivot(index="pH", columns="dataset", values="size")
            .fillna(0)
        )
        split_counts = split_counts.reindex(
            columns=["Training", "Validation", "Testing"], fill_value=0
        )
        fig, ax = plt.subplots(figsize=(8, 4))
        split_counts.plot(kind="bar", ax=ax, color=["#2878B5", "#F2A104", "#D64550"])
        ax.set_xlabel("pH")
        ax.set_ylabel("CV curves")
        ax.set_title("Training, validation and testing distribution")
        ax.tick_params(axis="x", rotation=0)
        ax.legend(title="Dataset", frameon=False)
        fig.tight_layout()
        st.pyplot(fig); plt.close(fig)

        st.markdown("### Mean CV comparison for pH 3–7")
        selected_temperature=st.selectbox("Select temperature for mean CV comparison",sorted(data.temperature_C.unique()),key="mean_cv_temperature")
        colours={3:"#e63946",4:"#f59e0b",5:"#10b981",6:"#168aad",7:"#7b2cbf"}; fig,ax=plt.subplots(figsize=(10.5,6.2)); plotted=[]
        for ph_value in [3,4,5,6,7]:
            subset=data[np.isclose(data.pH.astype(float),ph_value)&np.isclose(data.temperature_C.astype(float),selected_temperature)]
            result=mean_cv(subset)
            if result is None: continue
            voltage,current,current_sd=result; colour=colours[ph_value]
            ax.plot(voltage,current,color=colour,lw=2.5,label=f"pH {ph_value} (n={len(subset)})"); ax.fill_between(voltage,current-current_sd,current+current_sd,color=colour,alpha=.10); plotted.append(ph_value)
        ax.set(xlabel="Potential (V)",ylabel="Mean current (nA)",title=f"Mean CV Responses at {selected_temperature:g} °C"); ax.grid(alpha=.2); ax.legend(title="pH",frameon=False,ncol=3); fig.tight_layout(); st.pyplot(fig,use_container_width=True); plt.close(fig)
        missing=[ph for ph in [3,4,5,6,7] if ph not in plotted]
        if missing: st.warning("No curves at this temperature for pH: "+", ".join(map(str,missing)))

        st.markdown("### Raw CV comparison for pH 3–7")
        raw_control_1,raw_control_2=st.columns(2)
        with raw_control_1:
            raw_temperature=st.selectbox("Select temperature for raw CV comparison",sorted(data.temperature_C.unique()),key="raw_comparison_temperature")
        with raw_control_2:
            raw_replicate=st.selectbox("Select replicate",sorted(data.replicate.unique()),key="raw_comparison_replicate")
        fig,ax=plt.subplots(figsize=(10.5,6.2)); raw_plotted=[]
        for ph_value in [3,4,5,6,7]:
            subset=data[
                np.isclose(data.pH.astype(float),ph_value)
                & np.isclose(data.temperature_C.astype(float),raw_temperature)
                & (data.replicate.astype(int)==int(raw_replicate))
            ]
            if subset.empty: continue
            curve=subset.iloc[0]; ax.plot(curve.potential_V,np.asarray(curve.current_A)*1e9,color=colours[ph_value],lw=2.1,label=f"pH {ph_value}"); raw_plotted.append(ph_value)
        ax.set_xlabel("Potential (V)",fontweight="bold"); ax.set_ylabel("Raw current (nA)",fontweight="bold"); ax.set_title(f"Raw CV Curves at {raw_temperature:g} °C · Replicate {int(raw_replicate)}",fontweight="bold"); ax.grid(alpha=.2); ax.legend(title="pH",frameon=False,ncol=5); fig.tight_layout(); st.pyplot(fig,use_container_width=True); plt.close(fig)
        raw_missing=[ph for ph in [3,4,5,6,7] if ph not in raw_plotted]
        if raw_missing: st.warning("The selected replicate was unavailable for pH: "+", ".join(map(str,raw_missing)))

        chosen_ph = st.selectbox("Raw CV plot: pH", sorted(data.pH.unique()))
        fig, ax = plt.subplots(figsize=(8,5))
        for _, row in data[data.pH == chosen_ph].iterrows(): ax.plot(row.potential_V, row.current_A*1e9, alpha=.45, label=f"{row.temperature_C} °C")
        handles, labels = ax.get_legend_handles_labels(); unique = dict(zip(labels, handles)); ax.legend(unique.values(), unique.keys(), ncol=4)
        ax.set(xlabel="Potential (V)", ylabel="Current (nA)", title=f"Raw CV curves at pH {chosen_ph:g}"); st.pyplot(fig); plt.close(fig)

with feature_tab:
    if "split_data" not in st.session_state: st.info("Upload and read your Excel files first.")
    elif st.button("Extract CV features", type="primary"):
        with st.spinner("Extracting electrochemical and signal-shape features..."):
            features = pd.DataFrame(st.session_state.split_data.apply(extract_features, axis=1).tolist())
            features["dataset"] = st.session_state.split_data["dataset"].to_numpy()
            st.session_state.features = features
            for key in ["ranking", "selected", "results", "models", "predictions"]: st.session_state.pop(key, None)
    if "features" in st.session_state:
        features = st.session_state.features
        st.success(f"Extracted {len([c for c in features if c not in META+['dataset']])} features from {len(features)} curves.")
        st.dataframe(features, use_container_width=True, height=360)
        numerical_features=[c for c in features if c not in META+["dataset"] and pd.api.types.is_numeric_dtype(features[c])]
        view_feature = st.selectbox("Plot an extracted feature", numerical_features)
        sphere_data=features[["pH","temperature_C",view_feature]].replace([np.inf,-np.inf],np.nan).dropna()
        st.markdown("#### 3D spherical feature visualization")
        fig=plt.figure(figsize=(10,7)); ax=fig.add_subplot(111,projection="3d")
        spheres=ax.scatter(
            sphere_data["temperature_C"],sphere_data["pH"],sphere_data[view_feature],
            c=sphere_data["temperature_C"],cmap="turbo",s=85,marker="o",
            edgecolor="white",linewidth=.65,alpha=.92,depthshade=True,
        )
        ax.set_xlabel("Temperature (°C)",fontweight="bold",labelpad=10)
        ax.set_ylabel("pH",fontweight="bold",labelpad=10)
        ax.set_zlabel(view_feature.replace("_"," "),fontweight="bold",labelpad=10)
        ax.set_title(f"3D Spherical Distribution of {view_feature.replace('_',' ').title()}",fontweight="bold",pad=18)
        fig.colorbar(spheres,ax=ax,pad=.10,shrink=.70,label="Temperature (°C)")
        ax.view_init(elev=24,azim=42); fig.tight_layout(); st.pyplot(fig,use_container_width=True); plt.close(fig)

        st.markdown("### Calibration plot")
        calibration_temperature=st.selectbox("Select calibration temperature",sorted(features.temperature_C.unique()),key="calibration_temperature")
        calibration=features[np.isclose(features.temperature_C.astype(float),calibration_temperature)].groupby("pH")[view_feature].agg(["mean","std","count"]).reset_index().replace([np.inf,-np.inf],np.nan).dropna(subset=["mean"])
        if len(calibration)>=2:
            slope,intercept=np.polyfit(calibration.pH,calibration["mean"],1); fitted=slope*calibration.pH+intercept; ss_res=np.sum((calibration["mean"]-fitted)**2); ss_tot=np.sum((calibration["mean"]-calibration["mean"].mean())**2); calibration_r2=1-ss_res/ss_tot if ss_tot else np.nan
            fig,ax=plt.subplots(figsize=(8.5,5)); ax.errorbar(calibration.pH,calibration["mean"],yerr=calibration["std"].fillna(0),fmt="o",ms=9,capsize=5,color="#6f42c1",label="Mean ± SD"); ax.plot(calibration.pH,fitted,"--",color="#d33b83",lw=2,label="Linear calibration")
            ax.set(xlabel="pH",ylabel=view_feature.replace("_"," "),title=f"Calibration at {calibration_temperature:g} °C"); ax.grid(alpha=.2); ax.legend(frameon=False); fig.tight_layout(); st.pyplot(fig,use_container_width=True); plt.close(fig)
            c1,c2,c3=st.columns(3); c1.metric("Slope / sensitivity",f"{slope:.5g}"); c2.metric("Intercept",f"{intercept:.5g}"); c3.metric("R²",f"{calibration_r2:.4f}")
        else: st.warning("At least two pH levels are required at this temperature for calibration.")
        st.download_button("Download all extracted features", excel_bytes({"All_CV_Features": features}), "CV_extracted_features.xlsx")

with selection_tab:
    if "features" not in st.session_state: st.info("Complete feature extraction first.")
    elif st.button("Calculate feature ranking", type="primary"):
        with st.spinner("Running ANOVA, mutual information, Random Forest and permutation importance..."):
            train_f = st.session_state.features.query("dataset == 'Training'").reset_index(drop=True)
            ranking, clean = rank_features(train_f)
            st.session_state.ranking, st.session_state.clean_train = ranking, clean
    if "ranking" in st.session_state:
        ranking = st.session_state.ranking
        top_n = st.slider("Number of top features used for modelling", 1, len(ranking), min(15, len(ranking)))
        st.session_state.selected = ranking.Feature.head(top_n).tolist()
        st.dataframe(ranking, use_container_width=True, height=350)
        fig, ax = plt.subplots(figsize=(9, max(4, top_n*.32))); plot = ranking.head(top_n).sort_values("Combined_Score"); ax.barh(plot.Feature, plot.Combined_Score); ax.set_xlabel("Combined feature-selection score"); st.pyplot(fig)
        corr = st.session_state.clean_train[st.session_state.selected].corr()
        fig, ax = plt.subplots(figsize=(9,7)); sns.heatmap(corr, cmap="coolwarm", center=0, ax=ax); ax.set_title("Selected-feature correlation"); st.pyplot(fig)
        if len(st.session_state.selected) >= 2:
            x = StandardScaler().fit_transform(st.session_state.clean_train[st.session_state.selected]); pca = PCA(n_components=min(3, x.shape[1])).fit_transform(x)
            fig = plt.figure(figsize=(8,6)); ax = fig.add_subplot(111, projection="3d" if pca.shape[1] >= 3 else None)
            y = st.session_state.features.query("dataset == 'Training'").pH.to_numpy()
            if pca.shape[1] >= 3: sc=ax.scatter(pca[:,0],pca[:,1],pca[:,2],c=y,cmap="viridis"); ax.set_zlabel("PC3")
            else: sc=ax.scatter(pca[:,0],pca[:,1],c=y,cmap="viridis")
            ax.set_xlabel("PC1"); ax.set_ylabel("PC2"); fig.colorbar(sc, label="pH"); ax.set_title("PCA of selected features"); st.pyplot(fig)
        st.download_button("Download feature ranking", excel_bytes({"Feature_Ranking": ranking, "Selected_Features": pd.DataFrame({"Feature":st.session_state.selected})}), "CV_feature_selection.xlsx")

with training_tab:
    if "selected" not in st.session_state: st.info("Calculate a feature ranking and select features first.")
    else:
        selected_models = st.multiselect("Models to train", list(model_set(len(st.session_state.selected), 10)), default=list(model_set(len(st.session_state.selected), 10)))
        if st.button("Train and evaluate models", type="primary", disabled=not selected_models):
            f, selected = st.session_state.features, st.session_state.selected
            train, val, test = [f.query("dataset == @name").reset_index(drop=True) for name in ["Training","Validation","Testing"]]
            if val.empty or test.empty: st.error("Validation or test data is empty. Each pH-temperature group needs at least 3 curves.")
            else:
                xtr,ytr=train[selected].replace([np.inf,-np.inf],np.nan),train.pH; xv,yv=val[selected].replace([np.inf,-np.inf],np.nan),val.pH; xt,yt=test[selected].replace([np.inf,-np.inf],np.nan),test.pH
                models=model_set(len(selected),len(train)); rows=[]; fitted={}; predictions={}
                progress=st.progress(0)
                for idx,name in enumerate(selected_models):
                    try:
                        model=clone(models[name]).fit(xtr,ytr); pv=np.asarray(model.predict(xv)).reshape(-1); pt=np.asarray(model.predict(xt)).reshape(-1)
                        fitted[name]=model; predictions[name]={"validation":pv,"test":pt}
                        rows.append({"Model":name,"Validation_R2":r2_score(yv,pv),"Validation_RMSE":np.sqrt(mean_squared_error(yv,pv)),"Validation_MAE":mean_absolute_error(yv,pv),
                                     "Test_R2":r2_score(yt,pt),"Test_RMSE":np.sqrt(mean_squared_error(yt,pt)),"Test_MAE":mean_absolute_error(yt,pt)})
                    except Exception as exc: st.warning(f"{name} skipped: {exc}")
                    progress.progress((idx+1)/len(selected_models))
                results=pd.DataFrame(rows).sort_values(["Validation_RMSE","Validation_MAE"]).reset_index(drop=True); results.insert(0,"Rank",range(1,len(results)+1))
                st.session_state.update(results=results, models=fitted, predictions=predictions, test_frame=test)
        if "results" in st.session_state:
            st.success(f"Best validation model: {st.session_state.results.iloc[0].Model}")
            st.dataframe(st.session_state.results.style.format(precision=4), use_container_width=True)

with comparison_tab:
    if "results" not in st.session_state: st.info("Train the models first.")
    else:
        results, predictions, test = st.session_state.results, st.session_state.predictions, st.session_state.test_frame
        metric = st.selectbox("Comparison metric", ["Validation_RMSE","Validation_MAE","Validation_R2","Test_RMSE","Test_MAE","Test_R2"])
        fig, ax = plt.subplots(figsize=(9,5)); ordered=results.sort_values(metric,ascending="R2" not in metric); sns.barplot(data=ordered,x=metric,y="Model",ax=ax); st.pyplot(fig)
        name = st.selectbox("Detailed model plots", results.Model)
        actual=test.pH.to_numpy(); predicted=predictions[name]["test"]
        classes=np.sort(st.session_state.features.pH.unique())
        c1,c2=st.columns(2)
        with c1:
            fig,ax=plt.subplots(figsize=(6.2,5.7))
            points=ax.scatter(actual,predicted,c=test.temperature_C,cmap="plasma",s=72,edgecolor="white",linewidth=.7,label="Model-predicted pH")
            lo=min(actual.min(),predicted.min()); hi=max(actual.max(),predicted.max()); ideal_line,=ax.plot([lo,hi],[lo,hi],"k--",lw=2,label="Ideal prediction (actual = predicted)")
            ax.set_xlabel("Actual experimental pH",fontweight="bold"); ax.set_ylabel("Model-predicted pH",fontweight="bold"); ax.set_title(f"{name}\nActual vs Predicted pH",fontweight="bold",pad=12); ax.grid(alpha=.18)
            fig.colorbar(points,ax=ax,pad=.02,label="Temperature (°C)")
            ax.legend(handles=[points,ideal_line],loc="upper center",bbox_to_anchor=(.5,-.16),ncol=1,frameon=False,fontsize=8)
            fig.subplots_adjust(bottom=.27); st.pyplot(fig,use_container_width=True); plt.close(fig)
        with c2:
            fig,ax=plt.subplots(figsize=(5.5,5)); actual_classes=nearest(actual,classes); plotted=False
            for ph_class in classes:
                binary=(actual_classes==ph_class).astype(int)
                if binary.min()==binary.max(): continue
                score=-np.abs(predicted-ph_class); false_positive_rate,true_positive_rate,_=roc_curve(binary,score)
                ax.plot(false_positive_rate,true_positive_rate,lw=2,label=f"pH {ph_class:g} (AUC={auc(false_positive_rate,true_positive_rate):.2f})"); plotted=True
            ax.plot([0,1],[0,1],"k--",alpha=.6); ax.set_xlabel("False-positive rate"); ax.set_ylabel("True-positive rate"); ax.set_title(f"{name}: One-vs-Rest ROC")
            if plotted: ax.legend(fontsize=8,loc="upper center",bbox_to_anchor=(.5,-.16),ncol=2,frameon=False)
            ax.grid(alpha=.2); fig.subplots_adjust(bottom=.27); st.pyplot(fig); plt.close(fig)
        plot_r2=r2_score(actual,predicted); plot_rmse=np.sqrt(mean_squared_error(actual,predicted)); plot_mae=mean_absolute_error(actual,predicted)
        metric_1,metric_2,metric_3=st.columns(3); metric_1.metric("Test R²",f"{plot_r2:.4f}"); metric_2.metric("Test RMSE",f"{plot_rmse:.4f} pH"); metric_3.metric("Test MAE",f"{plot_mae:.4f} pH")
        st.caption("Actual and predicted values are listed in the table below; point labels are intentionally kept off the graph to prevent overlap.")
        classified=nearest(predicted,classes); cm=confusion_matrix(actual,classified,labels=classes)
        acc=accuracy_score(actual,classified); bal=balanced_accuracy_score(actual,classified)
        st.write(f"Nearest-pH classification accuracy: **{acc:.1%}** · Balanced accuracy: **{bal:.1%}**")
        fig,ax=plt.subplots(figsize=(6,5)); sns.heatmap(cm,annot=True,fmt="d",cmap="Blues",xticklabels=classes,yticklabels=classes,ax=ax); ax.set(xlabel="Predicted pH class",ylabel="Actual pH class",title=f"{name}: confusion matrix"); st.pyplot(fig)
        table=pd.DataFrame({"Sample ID":test.sample_id,"Temperature (°C)":test.temperature_C,"Actual experimental pH":actual,"Model-predicted pH":predicted,"Absolute error":abs(predicted-actual)})
        st.dataframe(table,use_container_width=True)
        st.download_button("Download model results",excel_bytes({"Model_Comparison":results,"Test_Predictions":table}),"CV_model_results.xlsx")
        buffer=io.BytesIO(); joblib.dump({"model":st.session_state.models[name],"selected_features":st.session_state.selected,"pH_classes":classes},buffer)
        st.download_button(f"Download trained {name}",buffer.getvalue(),f"{name.replace(' ','_')}_model.joblib")

with prediction_tab:
    if "results" not in st.session_state:
        st.info("Complete model training first. This page will then predict one randomly selected test CV.")
    else:
        results = st.session_state.results
        test = st.session_state.test_frame.reset_index(drop=True)
        predictions = st.session_state.predictions
        best_model = results.iloc[0]["Model"]
        st.subheader("🎯 Random test-sample prediction")
        st.caption("A test curve that was not used for model fitting is selected for an independent demonstration.")

        control_1, control_2 = st.columns([2, 1])
        with control_1:
            prediction_model = st.selectbox(
                "Prediction model", results["Model"].tolist(),
                index=results["Model"].tolist().index(best_model), key="random_prediction_model"
            )
        with control_2:
            if st.button("🎲 Select random test CV", type="primary", use_container_width=True):
                previous = st.session_state.get("random_test_index", -1)
                rng = np.random.default_rng()
                choices = [i for i in range(len(test)) if i != previous] or [0]
                st.session_state.random_test_index = int(rng.choice(choices))
        if "random_test_index" not in st.session_state or st.session_state.random_test_index >= len(test):
            st.session_state.random_test_index = 0

        random_index = st.session_state.random_test_index
        sample = test.iloc[random_index]
        predicted_ph = float(predictions[prediction_model]["test"][random_index])
        actual_ph = float(sample["pH"])
        absolute_error = abs(predicted_ph - actual_ph)

        p1, p2, p3, p4 = st.columns(4)
        with p1:
            st.markdown(f"""<div class="prediction-card"><div>AI PREDICTED pH</div>
            <div class="value">{predicted_ph:.3f}</div><small>{prediction_model}</small></div>""", unsafe_allow_html=True)
        p2.metric("🧪 Actual test pH", f"{actual_ph:g}")
        p3.metric("🌡️ Temperature", f"{float(sample['temperature_C']):g} °C")
        p4.metric("📏 Absolute error", f"{absolute_error:.3f}")

        raw_matches = st.session_state.split_data[
            st.session_state.split_data["sample_id"] == sample["sample_id"]
        ]
        raw = None
        left, right = st.columns([1.15, 1])
        with left:
            st.markdown("#### Raw test CV data plot")
            if not raw_matches.empty:
                raw = raw_matches.iloc[0]
                fig, ax = plt.subplots(figsize=(7, 4.6))
                ax.plot(raw["potential_V"], np.asarray(raw["current_A"]) * 1e9,
                        color="#7b2cbf", linewidth=2.3)
                ax.fill_between(raw["potential_V"], np.asarray(raw["current_A"]) * 1e9,
                                alpha=.15, color="#00a6c7")
                ax.set_xlabel("Potential (V)"); ax.set_ylabel("Current (nA)")
                ax.set_title(f"{sample['sample_id']} · Test measurement")
                ax.grid(alpha=.2); fig.tight_layout(); st.pyplot(fig)
                raw_test_csv = pd.DataFrame({
                    "potential_V": np.asarray(raw["potential_V"], dtype=float),
                    "current_A": np.asarray(raw["current_A"], dtype=float),
                    "current_nA": np.asarray(raw["current_A"], dtype=float) * 1e9,
                    "temperature_C": float(sample["temperature_C"]),
                    "actual_pH": actual_ph,
                    "sample_id": sample["sample_id"],
                })
                st.download_button(
                    "⬇️ Download selected raw test CV as CSV",
                    raw_test_csv.to_csv(index=False).encode("utf-8"),
                    f"{sample['sample_id']}_raw_test_CV.csv",
                    mime="text/csv",
                    use_container_width=True,
                    key=f"download_raw_test_{sample['sample_id']}",
                )
        with right:
            st.markdown("#### Actual and predicted pH")
            all_model_values = pd.DataFrame({
                "Model": ["Actual pH"] + results["Model"].tolist(),
                "pH": [actual_ph] + [float(predictions[m]["test"][random_index]) for m in results["Model"]],
                "Type": ["Experimental"] + ["AI prediction"] * len(results),
            })
            fig, ax = plt.subplots(figsize=(7, 4.6))
            colors = ["#0f766e"] + ["#7b2cbf"] * len(results)
            bars = ax.barh(all_model_values["Model"], all_model_values["pH"], color=colors)
            for bar, value in zip(bars, all_model_values["pH"]):
                ax.text(bar.get_width() + .03, bar.get_y() + bar.get_height()/2,
                        f"{value:.2f}", va="center", fontsize=8)
            ax.axvline(actual_ph, color="#0f766e", linestyle="--", linewidth=1.5)
            ax.set_xlabel("pH"); ax.set_title("One random test sample: all-model prediction")
            ax.invert_yaxis(); fig.tight_layout(); st.pyplot(fig)

        st.markdown("#### Features extracted from this test CV")
        selected = st.session_state.selected
        test_feature_table = pd.DataFrame({
            "Rank": np.arange(1, len(selected) + 1),
            "Selected feature": selected,
            "Extracted value": [sample[f] for f in selected],
        })
        feature_left, feature_right = st.columns([1, 1.2])
        with feature_left:
            st.dataframe(test_feature_table, use_container_width=True, hide_index=True, height=430)
        with feature_right:
            train_features = st.session_state.features.query("dataset == 'Training'")[selected]
            medians = train_features.replace([np.inf, -np.inf], np.nan).median()
            scales = train_features.replace([np.inf, -np.inf], np.nan).std().replace(0, 1)
            standardized = ((sample[selected].astype(float) - medians) / scales).clip(-5, 5)
            feature_plot = pd.DataFrame({"Feature": selected, "Standardized value": standardized.values})
            feature_plot = feature_plot.sort_values("Standardized value")
            fig, ax = plt.subplots(figsize=(7, max(4.5, len(selected) * .3)))
            bar_colors = ["#ef476f" if x < 0 else "#06a77d" for x in feature_plot["Standardized value"]]
            ax.barh(feature_plot["Feature"], feature_plot["Standardized value"], color=bar_colors)
            ax.axvline(0, color="#333", linewidth=1); ax.set_xlabel("Standardized feature value (z-score)")
            ax.set_title("Test-CV feature profile relative to training data"); fig.tight_layout(); st.pyplot(fig)

        st.markdown("### 📡 Live CV acquisition simulator")
        st.caption("This progressively streams the selected held-out CV. The same display can later receive live ADuCM355/API samples.")
        if st.button("▶ Play live CV stream",use_container_width=True,disabled=raw is None):
            voltage=np.asarray(raw["potential_V"],float); current=np.asarray(raw["current_A"],float)*1e9
            live_chart=st.empty(); live_status=st.empty(); live_progress=st.progress(0)
            frame_ends=np.linspace(max(5,len(voltage)//30),len(voltage),28,dtype=int)
            for frame_number,end in enumerate(frame_ends,1):
                completion=end/len(voltage); evolving=(actual_ph+predicted_ph)/2+(predicted_ph-(actual_ph+predicted_ph)/2)*completion
                fig,ax=plt.subplots(figsize=(10,4.5)); ax.plot(voltage[:end],current[:end],color="#6f42c1",lw=2.3); ax.fill_between(voltage[:end],current[:end],alpha=.13,color="#00a6c7")
                ax.set_xlim(voltage.min(),voltage.max()); pad=max(np.ptp(current)*.08,.1); ax.set_ylim(current.min()-pad,current.max()+pad); ax.set_xlabel("Potential (V)",fontweight="bold"); ax.set_ylabel("Current (nA)",fontweight="bold"); ax.set_title(f"Live electrochemical stream · {completion:.0%} acquired",fontweight="bold"); ax.grid(alpha=.18); fig.tight_layout()
                live_chart.pyplot(fig,clear_figure=True,use_container_width=True); plt.close(fig); live_status.markdown(f"<div class='section-card'><span class='status-dot'></span><b>Streaming</b> · points {end}/{len(voltage)} · evolving prediction <b>pH {evolving:.3f}</b></div>",unsafe_allow_html=True); live_progress.progress(frame_number/len(frame_ends)); time.sleep(.07)
            live_status.success(f"Acquisition complete · final predicted pH {predicted_ph:.3f} · actual pH {actual_ph:g}")

        st.markdown("### 🎬 Downloadable prediction animation")
        gif_key=f"{sample['sample_id']}::{prediction_model}"
        if st.button("Generate animated prediction",use_container_width=True,disabled=raw is None):
            with st.spinner("Animating CV acquisition and evolving prediction..."):
                st.session_state.prediction_video=prediction_gif(raw,predicted_ph,actual_ph); st.session_state.prediction_video_key=gif_key
        if st.session_state.get("prediction_video_key")==gif_key:
            st.image(st.session_state.prediction_video,caption="Machine-learning pH prediction animation",use_container_width=True)
            st.download_button("Download prediction GIF",st.session_state.prediction_video,"AI_pH_prediction.gif",mime="image/gif",use_container_width=True)

        report = pd.DataFrame({
            "sample_id": [sample["sample_id"]], "location": [location_name],
            "date_time": [local_now.strftime("%Y-%m-%d %H:%M:%S %Z")],
            "latitude": [latitude], "longitude": [longitude], "model": [prediction_model],
            "actual_pH": [actual_ph], "predicted_pH": [predicted_ph],
            "absolute_error": [absolute_error], "temperature_C": [sample["temperature_C"]],
        })
        st.download_button(
            "⬇️ Download this smart prediction report",
            excel_bytes({"Prediction_Summary": report, "Extracted_Test_Features": test_feature_table}),
            "random_test_pH_prediction.xlsx",
            use_container_width=True,
        )

        st.markdown("### Dynamically generated scientific reports")
        report_key=f"{prediction_model}::{location_name}::{len(test)}"; pdf_col,ppt_col=st.columns(2)
        with pdf_col:
            if st.button("Generate PDF report",type="primary",use_container_width=True):
                with st.spinner("Creating PDF..."):
                    pdf_report,_=create_pdf_and_powerpoint(location_name,local_now.strftime("%Y-%m-%d %H:%M:%S %Z"),results,test,predictions[prediction_model]["test"],prediction_model,include_powerpoint=False)
                    st.session_state.pdf_report=pdf_report; st.session_state.pdf_report_key=report_key
            if st.session_state.get("pdf_report_key")==report_key: st.download_button("Download PDF",st.session_state.pdf_report,"AI_TC_pH_Report.pdf",mime="application/pdf",use_container_width=True)
        with ppt_col:
            if st.button("Generate PowerPoint",type="primary",use_container_width=True):
                try:
                    with st.spinner("Creating PowerPoint..."):
                        _,ppt_report=create_pdf_and_powerpoint(location_name,local_now.strftime("%Y-%m-%d %H:%M:%S %Z"),results,test,predictions[prediction_model]["test"],prediction_model,include_powerpoint=True)
                        st.session_state.ppt_report=ppt_report; st.session_state.ppt_report_key=report_key
                except ImportError: st.error("Upload requirements.txt with the exact filename, then reboot Streamlit. It contains python-pptx.")
                except Exception as exc: st.error(f"PowerPoint generation failed: {exc}")
            if st.session_state.get("ppt_report_key")==report_key: st.download_button("Download PowerPoint",st.session_state.ppt_report,"AI_TC_pH_Report.pptx",mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",use_container_width=True)

with live_test_tab:
    st.subheader("📡 Live Testing Data")
    st.caption("Import a CSV saved by MIT App Inventor from your phone or OneDrive, then run the same trained feature-extraction and prediction workflow.")
    if "results" not in st.session_state or "models" not in st.session_state:
        st.info("Complete feature selection and model training first. The trained model is required for live-test prediction.")
    else:
        st.markdown("#### 1 · Access the MIT App CSV")
        live_source=st.radio("CSV source",["📱 Phone / OneDrive file picker","☁️ Paste OneDrive shared link"],horizontal=True,key="live_csv_source")
        live_file=None
        if live_source.startswith("📱"):
            live_file=st.file_uploader("Choose MIT App CSV",type=["csv"],key="live_csv_upload",help="On Android/iPhone, tap Browse files and choose OneDrive. If OneDrive is not shown, use the shared-link button instead.")
        else:
            st.info("In OneDrive, select the CSV → Share → Link settings → Anyone with the link can view → Copy link.")
            live_url=st.text_input("OneDrive CSV shared link",placeholder="https://1drv.ms/...")
            live_filename=st.text_input("CSV filename",value="MIT_live_test.csv")
            if st.button("☁️ Fetch CSV directly from OneDrive",type="primary",use_container_width=True):
                try:
                    downloaded=download_onedrive_file(live_url,live_filename)
                    st.session_state.live_onedrive_csv=downloaded.getvalue(); st.session_state.live_onedrive_name=live_filename
                    st.success("CSV downloaded from OneDrive.")
                except Exception as exc: st.error(f"OneDrive download failed: {exc}")
            if st.session_state.get("live_onedrive_csv"):
                live_file=io.BytesIO(st.session_state.live_onedrive_csv); live_file.name=st.session_state.get("live_onedrive_name","MIT_live_test.csv")

        if live_file is not None:
            try:
                live_csv=pd.read_csv(live_file)
                if live_csv.empty or len(live_csv.columns)<2: raise ValueError("The CSV is empty or has fewer than two columns.")
                st.success(f"Loaded {len(live_csv):,} rows from {live_file.name}")
                st.dataframe(live_csv.head(50),use_container_width=True,height=250)
                columns=list(live_csv.columns)
                def preferred(names,fallback):
                    lowered={str(c).strip().lower():c for c in columns}
                    for name in names:
                        if name in lowered: return columns.index(lowered[name])
                    return fallback
                c1,c2,c3=st.columns(3)
                with c1: potential_col=st.selectbox("Potential / voltage column",columns,index=preferred(["potential","potential_v","voltage","voltage_v","ewe/v"],0))
                with c2: current_col=st.selectbox("Current column",columns,index=preferred(["current","current_a","current_na","i/a"],min(1,len(columns)-1)))
                with c3:
                    temperature_candidates=["Enter one temperature"]+columns
                    temperature_col=st.selectbox("Temperature source",temperature_candidates,index=0)
                u1,u2,u3=st.columns(3)
                with u1: current_unit=st.selectbox("Current unit in CSV",["A","µA","nA"],index=0)
                with u2: manual_temperature=st.number_input("Test temperature (°C)",value=25.0,disabled=temperature_col!="Enter one temperature")
                with u3:
                    model_names=st.session_state.results["Model"].tolist(); live_model=st.selectbox("Prediction model",model_names,index=0,key="live_model")
                has_actual=st.checkbox("I know the experimental pH (optional validation)")
                actual_live=st.number_input("Actual experimental pH",min_value=0.0,max_value=14.0,value=5.0,step=.1,disabled=not has_actual)
                if st.button("🧠 Extract features and predict live pH",type="primary",use_container_width=True):
                    voltage=pd.to_numeric(live_csv[potential_col],errors="coerce")
                    current=pd.to_numeric(live_csv[current_col],errors="coerce")
                    valid=voltage.notna()&current.notna(); voltage=voltage[valid].to_numpy(float); current=current[valid].to_numpy(float)
                    if len(voltage)<10: raise ValueError("At least 10 valid potential-current rows are required.")
                    factor={"A":1.0,"µA":1e-6,"nA":1e-9}[current_unit]
                    if temperature_col=="Enter one temperature": temperature=float(manual_temperature)
                    else:
                        temp_values=pd.to_numeric(live_csv.loc[valid,temperature_col],errors="coerce"); temperature=float(temp_values.median())
                        if not np.isfinite(temperature): raise ValueError("The selected temperature column has no valid numeric values.")
                    raw=pd.Series({"sample_id":"MIT_Live_Test","pH":np.nan,"temperature_C":temperature,"replicate":1,"potential_V":voltage,"current_A":current*factor,"source_file":live_file.name})
                    feature_values=extract_features(raw); selected=st.session_state.selected
                    model_input=pd.DataFrame([feature_values])[selected].replace([np.inf,-np.inf],np.nan)
                    predicted=float(np.asarray(st.session_state.models[live_model].predict(model_input)).reshape(-1)[0])
                    st.session_state.live_result={"raw":raw,"features":feature_values,"selected":selected,"model":live_model,"predicted":predicted,"temperature":temperature,"actual":float(actual_live) if has_actual else None,"filename":live_file.name}
                    st.session_state.pop("live_prediction_gif",None); st.session_state.pop("live_reports",None)
            except Exception as exc: st.error(f"CSV processing failed: {exc}")

        if "live_result" in st.session_state:
            live=st.session_state.live_result; raw=live["raw"]; predicted=live["predicted"]; actual=live["actual"]
            st.divider(); st.markdown("#### 2 · Live-test prediction result")
            m1,m2,m3,m4=st.columns(4)
            with m1: st.markdown(f"<div class='prediction-card'><div>LIVE PREDICTED pH</div><div class='value'>{predicted:.3f}</div><small>{live['model']}</small></div>",unsafe_allow_html=True)
            m2.metric("🌡️ Temperature",f"{live['temperature']:.2f} °C")
            m3.metric("🧪 Nearest pH class",f"{nearest([predicted],sorted(st.session_state.features.pH.unique()))[0]:g}")
            m4.metric("📏 Absolute error","Not supplied" if actual is None else f"{abs(predicted-actual):.3f}")
            left,right=st.columns([1.15,1])
            with left:
                fig,ax=plt.subplots(figsize=(7,4.6)); current_na=np.asarray(raw["current_A"])*1e9; ax.plot(raw["potential_V"],current_na,color="#6f2dbd",lw=2.3); ax.fill_between(raw["potential_V"],current_na,color="#00a6c7",alpha=.14); ax.set(xlabel="Potential (V)",ylabel="Current (nA)",title="MIT App live-test CV"); ax.grid(alpha=.2); fig.tight_layout(); st.pyplot(fig)
            with right:
                feature_table=pd.DataFrame({"Selected feature":live["selected"],"Extracted value":[live["features"][name] for name in live["selected"]]})
                st.markdown("##### Features extracted from live test data"); st.dataframe(feature_table,use_container_width=True,hide_index=True,height=390)
            st.markdown("#### 3 · Video-like live acquisition")
            if st.button("▶ Play uploaded CV as a live stream",use_container_width=True):
                voltage=np.asarray(raw["potential_V"],float); current=np.asarray(raw["current_A"],float)*1e9; chart=st.empty(); status=st.empty(); progress=st.progress(0)
                for frame,end in enumerate(np.linspace(max(5,len(voltage)//30),len(voltage),28,dtype=int),1):
                    completion=end/len(voltage); evolving=predicted-.6*(1-completion); fig,ax=plt.subplots(figsize=(10,4.4)); ax.plot(voltage[:end],current[:end],color="#6f42c1",lw=2.3); ax.fill_between(voltage[:end],current[:end],color="#00a6c7",alpha=.13); ax.set_xlim(voltage.min(),voltage.max()); pad=max(np.ptp(current)*.08,.1); ax.set_ylim(current.min()-pad,current.max()+pad); ax.set(xlabel="Potential (V)",ylabel="Current (nA)",title=f"Live ADuCM355 acquisition · {completion:.0%}"); ax.grid(alpha=.18); fig.tight_layout(); chart.pyplot(fig,clear_figure=True,use_container_width=True); plt.close(fig); status.markdown(f"<div class='section-card'><span class='status-dot'></span><b>Streaming</b> · {end}/{len(voltage)} points · evolving prediction <b>pH {evolving:.3f}</b></div>",unsafe_allow_html=True); progress.progress(frame/28); time.sleep(.06)
                status.success(f"Acquisition complete · predicted pH {predicted:.3f}")
            if st.button("🎬 Generate live prediction GIF",use_container_width=True):
                with st.spinner("Generating animation..."): st.session_state.live_prediction_gif=prediction_gif(raw,predicted,actual)
            if st.session_state.get("live_prediction_gif"):
                st.image(st.session_state.live_prediction_gif,caption="Uploaded MIT App CSV replay and evolving prediction",use_container_width=True); st.download_button("Download live prediction GIF",st.session_state.live_prediction_gif,"MIT_live_pH_prediction.gif",mime="image/gif",use_container_width=True)
            summary=pd.DataFrame({"source_file":[live["filename"]],"location":[location_name],"date_time":[local_now.strftime("%Y-%m-%d %H:%M:%S %Z")],"model":[live["model"]],"temperature_C":[live["temperature"]],"actual_pH":[actual],"predicted_pH":[predicted],"absolute_error":[None if actual is None else abs(predicted-actual)]})
            st.download_button("⬇️ Download live-test Excel results",excel_bytes({"Live_Prediction":summary,"Extracted_Features":feature_table}),"MIT_live_test_results.xlsx",use_container_width=True)
            if st.button("📄 Generate live-test PDF and PowerPoint",use_container_width=True):
                try:
                    with st.spinner("Creating reports..."): st.session_state.live_reports=create_live_reports(location_name,local_now.strftime("%Y-%m-%d %H:%M:%S %Z"),live["model"],predicted,live["temperature"],raw,feature_table,actual)
                except ImportError: st.error("python-pptx is missing. Add the supplied requirements.txt to GitHub and reboot the app.")
                except Exception as exc: st.error(f"Report generation failed: {exc}")
            if st.session_state.get("live_reports"):
                pdf_live,ppt_live=st.session_state.live_reports; d1,d2=st.columns(2); d1.download_button("Download live-test PDF",pdf_live,"MIT_live_test_report.pdf",mime="application/pdf",use_container_width=True); d2.download_button("Download live-test PowerPoint",ppt_live,"MIT_live_test_report.pptx",mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",use_container_width=True)

st.sidebar.header("Workflow status")
for label,key in [("Data uploaded","cv_data"),("Features extracted","features"),("Features selected","selected"),("Models trained","results")]:
    st.sidebar.write(("✅" if key in st.session_state else "○")+" "+label)
st.sidebar.info("Keep this browser tab open while processing. Results are retained while the Streamlit session remains active.")
st.sidebar.markdown(
    """
    <div style="margin-top:1.2rem;padding:1rem .75rem;border-top:1px solid rgba(255,255,255,.28);text-align:center;line-height:1.55">
      <div style="font-size:.78rem;opacity:.78;letter-spacing:.06em">RESEARCH APPLICATION</div>
      <div style="font-size:1rem;font-weight:750">Developed by Akshaya AV</div>
      <div style="font-size:.76rem;opacity:.75">© 2026 · AI-TC pH Sensor App</div>
    </div>
    """,
    unsafe_allow_html=True,
)
